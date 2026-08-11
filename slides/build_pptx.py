#!/usr/bin/env python3
"""Bikin presentasi.pptx dari presentasi.md.

Hasilnya di-import ke Google Slides. Teksnya tetap berupa teks (bukan
gambar), jadi masih bisa diedit dan di-restyle di sana.

    pip install python-pptx
    python3 slides/build_pptx.py

Skrip ini cuma paham subset Markdown yang dipakai di presentasi.md:
heading, paragraf, bullet, list bernomor, blockquote, tabel, dan code block.
"""

import math
import re
import sys
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ----------------------------------------------------------------- tampilan

FONT_BODY = "Arial"
FONT_MONO = "Courier New"

INK = RGBColor(0x1B, 0x25, 0x33)  # judul, teks utama
BODY = RGBColor(0x37, 0x44, 0x51)  # paragraf
MUTED = RGBColor(0x6B, 0x77, 0x84)  # nomor slide, keterangan
ACCENT = RGBColor(0x00, 0xAD, 0xD8)  # biru Go
CODE_BG = RGBColor(0xF3, 0xF6, 0xF8)
CODE_INK = RGBColor(0x1B, 0x27, 0x33)
TABLE_HEAD = RGBColor(0xE8, 0xF6, 0xFA)
RULE = RGBColor(0xDD, 0xE3, 0xE8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W, SLIDE_H = 13.333, 7.5
MARGIN = 0.85
CONTENT_W = SLIDE_W - 2 * MARGIN
BODY_TOP = 1.72
BODY_LIMIT = 6.80  # batas bawah sebelum nomor slide

SZ_BODY, SZ_CODE, SZ_TABLE, SZ_H3 = 17.0, 13.0, 14.0, 23.0

# Jarak antar blok, dalam inci. Blok yang punya "kotak" sendiri (code dan
# tabel) butuh napas lebih besar, kalau nggak barisnya kelihatan nempel
# ke paragraf di atasnya.
GAP_DEFAULT = 0.17
GAP_BOXED = 0.24
GAP_LIST = 0.20
GAP_BEFORE_TABLE = 0.12
GAP_LEAD_BEFORE = 0.14
GAP_LEAD_AFTER = 0.03

# Perkiraan lebar rata-rata satu karakter, relatif terhadap ukuran font.
# Dipakai buat menebak berapa baris sebuah paragraf akan makan tempat,
# karena python-pptx nggak bisa mengukur teks beneran.
CHAR_W_BODY = 0.50
CHAR_W_MONO = 0.601


@dataclass
class Block:
    kind: str
    text: str = ""
    items: list = field(default_factory=list)
    rows: list = field(default_factory=list)
    lang: str = ""


# ------------------------------------------------------------------ parsing

FENCE = "```"


def split_slides(md: str) -> list[list[str]]:
    """Pisah per slide di garis '---', tapi abaikan yang ada di dalam code block."""
    slides, cur, in_code = [], [], False
    for line in md.splitlines():
        if line.strip().startswith(FENCE):
            in_code = not in_code
        if line.strip() == "---" and not in_code:
            slides.append(cur)
            cur = []
        else:
            cur.append(line)
    slides.append(cur)
    return [s for s in slides if any(l.strip() for l in s)]


def is_special(s: str) -> bool:
    return (
        s.startswith(("#", ">", "- ", "|", FENCE))
        or re.match(r"^\d+\.\s", s) is not None
    )


def dedent(lines: list[str]) -> list[str]:
    real = [l for l in lines if l.strip()]
    if not real:
        return lines
    pad = min(len(l) - len(l.lstrip()) for l in real)
    return [l[pad:] if len(l) >= pad else l for l in lines]


def parse_slide(lines: list[str]) -> list[Block]:
    blocks: list[Block] = []
    i, n = 0, len(lines)
    while i < n:
        s = lines[i].strip()
        if not s:
            i += 1
            continue

        if s.startswith(FENCE):
            lang = s[len(FENCE) :].strip()
            i += 1
            code = []
            while i < n and not lines[i].strip().startswith(FENCE):
                code.append(lines[i])
                i += 1
            i += 1
            blocks.append(Block("code", lang=lang, items=dedent(code)))
            continue

        for marker, kind in (("### ", "h3"), ("## ", "h2"), ("# ", "h1")):
            if s.startswith(marker):
                blocks.append(Block(kind, s[len(marker) :].strip()))
                i += 1
                break
        else:
            if s.startswith("|"):
                rows = []
                while i < n and lines[i].strip().startswith("|"):
                    raw = lines[i].strip()
                    i += 1
                    if set(raw) <= set("|-: "):  # baris pemisah header
                        continue
                    rows.append([c.strip() for c in raw.strip("|").split("|")])
                blocks.append(Block("table", rows=rows))
            elif s.startswith("> "):
                blocks.append(Block("quote", s[2:].strip()))
                i += 1
            elif s.startswith("- "):
                items = []
                while i < n and lines[i].strip().startswith("- "):
                    items.append(lines[i].strip()[2:])
                    i += 1
                blocks.append(Block("bullets", items=items))
            elif re.match(r"^\d+\.\s", s):
                items = []
                while i < n and re.match(r"^\d+\.\s", lines[i].strip()):
                    items.append(lines[i].strip())
                    i += 1
                blocks.append(Block("numbers", items=items))
            else:
                para = [s]
                i += 1
                while i < n and lines[i].strip() and not is_special(lines[i].strip()):
                    para.append(lines[i].strip())
                    i += 1
                # Baris pertama yang seluruhnya bold diperlakukan sebagai
                # judul kecil. Kalau digabung jadi satu paragraf, isinya
                # nyambung terus dan slide-nya kelihatan seperti blok teks.
                lead = re.match(r"^\*\*([^*]+)\*\*$", para[0])
                if lead and len(para) > 1:
                    blocks.append(Block("lead", lead.group(1)))
                    blocks.append(Block("para", " ".join(para[1:])))
                else:
                    blocks.append(Block("para", " ".join(para)))
    return blocks


# ------------------------------------------------------- perkiraan tinggi

def plain(text: str) -> str:
    """Buang penanda markdown supaya panjangnya mendekati hasil render."""
    return re.sub(r"[*`]", "", text)


def wrapped_lines(text: str, width_in: float, size_pt: float, char_w: float) -> int:
    per_line = max(1, int((width_in * 72) / (size_pt * char_w)))
    return max(1, math.ceil(len(plain(text)) / per_line))


def block_height(b: Block, scale: float) -> float:
    """Perkiraan tinggi satu blok dalam inci."""
    if b.kind == "code":
        size = SZ_CODE * scale
        rows = max(1, len(b.items))
        return (rows * size * 1.32 + 20) / 72

    if b.kind == "table":
        size = SZ_TABLE * scale
        return len(b.rows) * (size * 2.05) / 72

    if b.kind == "h3":
        return (SZ_H3 * scale * 1.35) / 72

    size = SZ_BODY * scale
    if b.kind in ("bullets", "numbers"):
        total = sum(
            wrapped_lines(t, CONTENT_W - 0.35, size, CHAR_W_BODY) for t in b.items
        )
        return (total * size * 1.30 + len(b.items) * size * 0.34) / 72

    return (wrapped_lines(b.text, CONTENT_W, size, CHAR_W_BODY) * size * 1.30) / 72


def gap_before(b: Block) -> float:
    if b.kind == "table":
        return GAP_BEFORE_TABLE
    if b.kind == "lead":
        return GAP_LEAD_BEFORE  # pisahkan dari entri sebelumnya
    return 0.0


def gap_after(b: Block) -> float:
    if b.kind in ("code", "table"):
        return GAP_BOXED
    if b.kind in ("bullets", "numbers"):
        return GAP_LIST
    if b.kind == "lead":
        return GAP_LEAD_AFTER  # rapat, biar nempel ke penjelasannya
    return GAP_DEFAULT


def content_height(blocks: list[Block], scale: float) -> float:
    body = [b for b in blocks if b.kind not in ("h1", "h2")]
    if not body:
        return 0.0
    total = sum(block_height(b, scale) + gap_before(b) * scale for b in body)
    total += sum(gap_after(b) * scale for b in body[:-1])
    return total


def fit_scale(blocks: list[Block]) -> float:
    """Kecilkan font kalau isinya kepanjangan, tapi jangan sampai nggak kebaca."""
    room = BODY_LIMIT - BODY_TOP
    scale = 1.0
    for _ in range(14):
        h = content_height(blocks, scale)
        if h <= room or scale <= 0.62:
            break
        scale = max(0.62, scale * math.sqrt(room / h))
    return scale


# ------------------------------------------------------------- render teks

TOKEN = re.compile(r"(\*\*.+?\*\*|`[^`]+`|\*[^*\s][^*]*?\*)")


def add_runs(para, text, size, color, bold=False, mono=False):
    parts = [p for p in TOKEN.split(text) if p]
    if not parts:
        parts = [""]
    for tok in parts:
        run = para.add_run()
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = FONT_MONO if mono else FONT_BODY
        if tok.startswith("**") and tok.endswith("**") and len(tok) > 4:
            run.text = tok[2:-2]
            run.font.bold = True
        elif tok.startswith("`") and tok.endswith("`") and len(tok) > 2:
            run.text = tok[1:-1]
            run.font.name = FONT_MONO
            # Courier New kelihatan lebih besar dari Arial di ukuran yang sama,
            # jadi dikecilkan sedikit supaya sejajar dengan teks sekitarnya.
            run.font.size = Pt(size * 0.90)
        elif tok.startswith("*") and tok.endswith("*") and len(tok) > 2:
            run.text = tok[1:-1]
            run.font.italic = True
        else:
            run.text = tok


def textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return box, tf


# ------------------------------------------------------------ render slide

def render_title_slide(slide, blocks):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(MARGIN), Inches(2.62), Inches(1.5), Inches(0.075)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    bar.shadow.inherit = False

    title = next((b for b in blocks if b.kind == "h1"), None)
    _, tf = textbox(slide, MARGIN, 1.55, CONTENT_W, 1.0)
    p = tf.paragraphs[0]
    add_runs(p, title.text if title else "", 42, INK, bold=True)

    y = 3.0
    for b in blocks:
        if b.kind == "h1":
            continue
        if b.kind == "h3":
            _, tf = textbox(slide, MARGIN, y, CONTENT_W, 0.5)
            add_runs(tf.paragraphs[0], b.text, 22, ACCENT, bold=True, mono=True)
            y += 0.62
        elif b.kind == "code":
            y = render_code(slide, b, y, 12.0) + 0.2
        elif b.kind == "para":
            _, tf = textbox(slide, MARGIN, y, CONTENT_W, 0.42)
            add_runs(tf.paragraphs[0], b.text, 17, BODY)
            y += 0.44


def render_code(slide, block, y, size):
    lines = block.items or [""]
    longest = max((len(l) for l in lines), default=1)

    # Kecilkan font kalau ada baris yang lebih lebar dari area konten,
    # supaya code block nggak kepotong di kanan.
    max_w = (CONTENT_W - 0.34) * 72
    if longest * size * CHAR_W_MONO > max_w:
        size = max(8.0, max_w / (longest * CHAR_W_MONO))

    height = (len(lines) * size * 1.32 + 20) / 72
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(MARGIN), Inches(y), Inches(CONTENT_W), Inches(height)
    )
    shape.adjustments[0] = 0.04
    shape.fill.solid()
    shape.fill.fore_color.rgb = CODE_BG
    shape.line.color.rgb = RULE
    shape.line.width = Pt(0.75)
    shape.shadow.inherit = False

    tf = shape.text_frame
    tf.word_wrap = False
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = Inches(0.17)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.09)
    tf.margin_bottom = Inches(0.09)

    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.06
        run = p.add_run()
        run.text = line if line else " "
        run.font.name = FONT_MONO
        run.font.size = Pt(size)
        run.font.color.rgb = CODE_INK
    return y + height


def render_table(slide, block, y, size):
    rows, cols = len(block.rows), len(block.rows[0])
    row_h = (size * 2.05) / 72
    height = rows * row_h
    shape = slide.shapes.add_table(
        rows, cols, Inches(MARGIN), Inches(y), Inches(CONTENT_W), Inches(height)
    )
    table = shape.table
    table.first_row = True
    table.horz_banding = False

    for r, row in enumerate(block.rows):
        table.rows[r].height = Inches(row_h)
        for c, cell_text in enumerate(row):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_HEAD if r == 0 else WHITE
            cell.margin_left = cell.margin_right = Inches(0.13)
            cell.margin_top = cell.margin_bottom = Inches(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            add_runs(p, cell_text, size, INK if r == 0 else BODY, bold=(r == 0))
    return y + height


def render_content_slide(slide, blocks, number):
    head = next((b for b in blocks if b.kind == "h2"), None)
    if head:
        _, tf = textbox(slide, MARGIN, 0.60, CONTENT_W, 0.72)
        add_runs(tf.paragraphs[0], head.text, 29, INK, bold=True)
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(MARGIN), Inches(1.35), Inches(1.05), Inches(0.055)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background()
        bar.shadow.inherit = False

    body = [b for b in blocks if b.kind not in ("h1", "h2")]
    scale = fit_scale(blocks)
    y = BODY_TOP

    for b in body:
        y += gap_before(b) * scale
        if b.kind == "code":
            y = render_code(slide, b, y, SZ_CODE * scale)
        elif b.kind == "table":
            y = render_table(slide, b, y, SZ_TABLE * scale)
        elif b.kind == "h3":
            h = block_height(b, scale)
            _, tf = textbox(slide, MARGIN, y, CONTENT_W, h)
            add_runs(tf.paragraphs[0], b.text, SZ_H3 * scale, ACCENT, bold=True, mono=True)
            y += h
        elif b.kind in ("bullets", "numbers"):
            h = block_height(b, scale)
            _, tf = textbox(slide, MARGIN, y, CONTENT_W, h)
            for idx, item in enumerate(b.items):
                p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
                p.line_spacing = 1.28
                p.space_after = Pt(SZ_BODY * scale * 0.34)
                if b.kind == "bullets":
                    bullet = p.add_run()
                    bullet.text = "•   "
                    bullet.font.size = Pt(SZ_BODY * scale)
                    bullet.font.color.rgb = ACCENT
                    bullet.font.bold = True
                    bullet.font.name = FONT_BODY
                    add_runs(p, item, SZ_BODY * scale, BODY)
                else:
                    num, _, rest = item.partition(" ")
                    lead = p.add_run()
                    lead.text = num + "  "
                    lead.font.size = Pt(SZ_BODY * scale)
                    lead.font.color.rgb = ACCENT
                    lead.font.bold = True
                    lead.font.name = FONT_BODY
                    add_runs(p, rest, SZ_BODY * scale, BODY)
            y += h
        elif b.kind == "quote":
            h = block_height(b, scale)
            bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(MARGIN), Inches(y), Inches(0.045), Inches(h)
            )
            bar.fill.solid()
            bar.fill.fore_color.rgb = ACCENT
            bar.line.fill.background()
            bar.shadow.inherit = False
            _, tf = textbox(slide, MARGIN + 0.22, y, CONTENT_W - 0.22, h)
            p = tf.paragraphs[0]
            p.line_spacing = 1.30
            add_runs(p, b.text, SZ_BODY * scale, INK)
            y += h
        else:
            h = block_height(b, scale)
            _, tf = textbox(slide, MARGIN, y, CONTENT_W, h)
            p = tf.paragraphs[0]
            p.line_spacing = 1.30
            is_lead = b.kind == "lead"
            add_runs(p, b.text, SZ_BODY * scale, INK if is_lead else BODY, bold=is_lead)
            y += h
        y += gap_after(b) * scale

    _, tf = textbox(slide, SLIDE_W - MARGIN - 1.0, SLIDE_H - 0.62, 1.0, 0.3)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = str(number)
    run.font.size = Pt(11)
    run.font.color.rgb = MUTED
    run.font.name = FONT_BODY
    return scale


# ---------------------------------------------------------------------- main

def main() -> int:
    here = Path(__file__).resolve().parent
    src = here / "presentasi.md"
    out = here / "presentasi.pptx"

    decks = split_slides(src.read_text(encoding="utf-8"))
    parsed = [parse_slide(s) for s in decks]

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(SLIDE_W), Inches(SLIDE_H)
    blank = prs.slide_layouts[6]

    shrunk = []
    for idx, blocks in enumerate(parsed, start=1):
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = WHITE
        if any(b.kind == "h1" for b in blocks):
            render_title_slide(slide, blocks)
        else:
            scale = render_content_slide(slide, blocks, idx)
            if scale < 0.999:
                title = next((b.text for b in blocks if b.kind == "h2"), "?")
                shrunk.append((idx, title, scale))

    prs.save(out)

    print(f"{len(parsed)} slide -> {out}")
    if shrunk:
        print("\nFont dikecilkan supaya muat:")
        for idx, title, scale in shrunk:
            print(f"  slide {idx:>2}  {scale:.0%}  {title}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
